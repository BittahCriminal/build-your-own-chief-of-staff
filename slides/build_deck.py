#!/usr/bin/env python3
"""Build the follow-along talk deck. Non-technical. Speaker notes cite Notion."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

# Black + gold (49ers palette, no marks). Not a licensed theme.
BLACK = RGBColor(0x00, 0x00, 0x00)
CARD = RGBColor(0x14, 0x14, 0x14)
CREAM = RGBColor(0xF4, 0xEB, 0xD0)
GOLD = RGBColor(0xB3, 0x99, 0x5D)
MUTED = RGBColor(0xB8, 0xA8, 0x88)
NAVY = BLACK
NAVY2 = CARD

W, H = Inches(13.333), Inches(7.5)


def set_run(run, size=20, bold=False, color=CREAM, font="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    rPr = run._r.get_or_add_rPr()
    latin = rPr.find(qn("a:latin"))
    if latin is None:
        from pptx.oxml import parse_xml
        rPr.append(parse_xml(f'<a:latin xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" typeface="{font}"/>'))
    else:
        latin.set("typeface", font)


def add_rect(slide, l, t, w, h, fill):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    return sh


def textbox(slide, l, t, w, h, text, size=20, bold=False, color=CREAM, align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color, font=font)
    return box


def bullets(slide, l, t, w, h, items, size=22, color=CREAM, spacing=10):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(spacing)
        run = p.add_run()
        run.text = "•  " + item
        set_run(run, size=size, color=color)
    return box


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text.strip()


def footer(slide, n, total):
    textbox(
        slide,
        Inches(0.7),
        Inches(7.1),
        Inches(10),
        Inches(0.3),
        "Build your own chief of staff  ·  process, not a vendor",
        size=11,
        color=MUTED,
    )
    textbox(
        slide,
        Inches(11.6),
        Inches(7.1),
        Inches(1.2),
        Inches(0.3),
        f"{n}  /  {total}",
        size=11,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )


def gold_bar(slide):
    add_rect(slide, Inches(0), Inches(0), Inches(0.18), H, GOLD)
    add_rect(slide, Inches(0), Inches(7.42), W, Inches(0.08), GOLD)


def blank(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, W, H, NAVY)
    gold_bar(s)
    return s


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    TOTAL = 24

    def fin(s, n, note):
        footer(s, n, TOTAL)
        notes(s, note)

    # 1 Title
    s = blank(prs)
    textbox(s, Inches(0.8), Inches(1.7), Inches(11.5), Inches(0.4),
            "A TALK FOR OPERATORS, NOT DEVELOPERS", size=14, bold=True, color=GOLD)
    textbox(s, Inches(0.8), Inches(2.2), Inches(11.5), Inches(1.6),
            "Build your own\nchief of staff agent", size=44, bold=True, color=CREAM)
    textbox(s, Inches(0.8), Inches(4.5), Inches(11.5), Inches(1.2),
            "Teach the process. Not a vendor.\nLeave knowing how to build a CoS agent — and how to automate\nanything that is repeatable and verifiable.",
            size=20, color=MUTED)
    fin(s, 1, """
Open by saying this is not a product demo. No one is leaving with a login, a marketplace skill, or a reason to switch from Copilot to Claude.

They are leaving with a process they can run Monday in whatever tool they already have.

The Chief of Staff agent is the worked example. Briefs, drafts, open loops, meeting prep. The real prize is the two tests they will use on expense reports, hiring screens, and customer research next month.

Cite the source of truth if asked: the author's Substack research table in Notion — https://app.notion.com/p/36e059b703e180d3a962d862c9e380c5 — 879 rows of operator essays. We synthesized from those rows, not from a generic web roundup.
""")

    # 2 Walk away with
    s = blank(prs)
    textbox(s, Inches(0.8), Inches(0.45), Inches(11), Inches(0.8),
            "What you walk away with", size=32, bold=True)
    bullets(s, Inches(0.8), Inches(1.6), Inches(11.5), Inches(3.2), [
        "How to stand up a personal Chief of Staff without writing code.",
        "How to decide what should be an agent — and what must stay human.",
        "A folder of markdown files you paste into Copilot, Cursor, Claude, ChatGPT, or Gemini.",
    ], size=24, spacing=18)
    textbox(s, Inches(0.8), Inches(5.3), Inches(11.5), Inches(1.2),
            "The model is interchangeable. Your process is not.",
            size=22, bold=True, color=GOLD)
    fin(s, 2, """
Pause on the third bullet. These are ordinary markdown files. If their IT department only allows Copilot, they are fine. If they live in ChatGPT, they are fine.

The context files they fill in privately — who they are, this quarter's priorities, how they write, where truth lives — are the operating system. Swap the vendor tomorrow and nothing important moves. START-HERE.md is how the folder lands on a PC. Do not walk the file tree here.

That claim is from the context-files row in the table (Notion: https://app.notion.com/p/3bb059b703e18152aea0d83b502fa319) and from the delegation kit's memory scaffold (https://app.notion.com/p/36f059b703e18192a6d8f67fbe74b756).
""")

    # 3 Not a vendor
    s = blank(prs)
    textbox(s, Inches(0.8), Inches(0.45), Inches(11), Inches(0.8),
            "This is not a vendor talk", size=32, bold=True)
    add_rect(s, Inches(0.8), Inches(1.6), Inches(5.5), Inches(4.4), NAVY2)
    add_rect(s, Inches(6.9), Inches(1.6), Inches(5.5), Inches(4.4), NAVY2)
    textbox(s, Inches(1.05), Inches(1.85), Inches(5), Inches(0.5), "WE WILL NOT", size=14, bold=True, color=GOLD)
    bullets(s, Inches(1.05), Inches(2.4), Inches(5), Inches(3.2), [
        "Pick a model for you.",
        "Install a plugin.",
        "Promise it replaces a human CoS.",
        "Let it send anything.",
    ], size=18, spacing=12)
    textbox(s, Inches(7.15), Inches(1.85), Inches(5), Inches(0.5), "WE WILL", size=14, bold=True, color=GOLD)
    bullets(s, Inches(7.15), Inches(2.4), Inches(5), Inches(3.2), [
        "Write the job down.",
        "Put you in files you own.",
        "Check the draft against a source.",
        "Fix the file when it is wrong.",
    ], size=18, spacing=12)
    fin(s, 3, """
Say the refusal out loud. 'It replaces your Chief of Staff' is marketing. A human CoS does two kinds of work. Operational: reconstruct context, draft the update, prep the meeting, keep loops from disappearing. Judgment: when to push, how a room will land, which red line you will not cross.

We automate the first. We keep the second. That split is the whole talk, and it is the load-bearing claim in the Delegation Kit row: https://app.notion.com/p/36f059b703e18192a6d8f67fbe74b756
""")

    # 4 Two tests
    s = blank(prs)
    textbox(s, Inches(0.8), Inches(0.45), Inches(11), Inches(0.8),
            "The whole course is two tests", size=32, bold=True)
    add_rect(s, Inches(0.8), Inches(1.7), Inches(5.5), Inches(4.3), NAVY2)
    add_rect(s, Inches(6.9), Inches(1.7), Inches(5.5), Inches(4.3), NAVY2)
    textbox(s, Inches(1.1), Inches(2.0), Inches(5), Inches(0.6), "1. Repeatable", size=26, bold=True, color=GOLD)
    textbox(s, Inches(1.1), Inches(2.8), Inches(5), Inches(2.6),
            "You would do the same steps next week.\nSame kinds of inputs. Same shape of output.\n\nSpecial cases every time = judgment.\nKeep it.",
            size=18, color=CREAM)
    textbox(s, Inches(7.2), Inches(2.0), Inches(5), Inches(0.6), "2. Verifiable", size=26, bold=True, color=GOLD)
    textbox(s, Inches(7.2), Inches(2.8), Inches(5), Inches(2.6),
            "A careful human can check the output\nagainst a source, a template, or a\nyes/no rule.\n\n“Does this sound good?” is not a check.",
            size=18, color=CREAM)
    fin(s, 4, """
Stay here. This is the slide they should photograph.

Repeatable means the steps do not change because the personalities in the room changed. Verifiable means someone who was not in the model's head can open a ticket, a note, or a template and say pass or fail.

Fail either test → do not automate. That is not caution. That is the curriculum.

The checkability argument in the table: if checking an answer costs as much as making it, extra attempts just grow the pile. Notion: https://app.notion.com/p/399059b703e18107b20dfe6bb5c32626  (One-Minute Test / agent-shaped work)

And: if you cannot name what would make you say 'not yet,' you have a vibe, not a job. Notion: https://app.notion.com/p/36f059b703e18142b351f70732b09c29
""")

    # 5 Repeatable
    s = blank(prs)
    textbox(s, Inches(0.8), Inches(0.45), Inches(11), Inches(0.8),
            "Repeatable = same steps next week", size=32, bold=True)
    bullets(s, Inches(0.8), Inches(1.6), Inches(11.5), Inches(4.8), [
        "Monday status to your manager: same headings, new evidence. Yes.",
        "Meeting recap: decisions, owners, dates. Yes.",
        "Which of two candidates to hire: no. The facts can be gathered; the call cannot.",
        "“Handle my email” is usually seventeen jobs you have never named. Split it.",
    ], size=22, spacing=16)
    fin(s, 5, """
Give them the email example from the Delegation Kit. People say 'handle my email' and mean: the Monday status, the recap, the ask, the bad-news note, the FYI they should have sent yesterday. Those are different jobs with different checks.

The first-agent-job guide is blunt: start with one repeated problem and three real examples, not an agent idea. Notion: https://app.notion.com/p/3bb059b703e1817b9123ff209fcc5d9d

If they only have a vibe, they do not have a first job yet. That is allowed.
""")

    # 6 Verifiable
    s = blank(prs)
    textbox(s, Inches(0.8), Inches(0.45), Inches(11), Inches(0.8),
            "Verifiable = a human can check it", size=32, bold=True)
    add_rect(s, Inches(0.8), Inches(1.6), Inches(11.5), Inches(4.6), NAVY2)
    textbox(s, Inches(1.15), Inches(1.9), Inches(10.8), Inches(0.5),
            "A CHECK looks like", size=14, bold=True, color=GOLD)
    bullets(s, Inches(1.15), Inches(2.5), Inches(10.8), Inches(3.3), [
        "Every Done line cites a ticket, PR, or note I can open.",
        "The recap names an owner only when the notes name one. Otherwise: unassigned.",
        "The brief lists sources it could not see, instead of pretending it read them.",
        "Length: one screen. Hedging language: none. Recommendation: present.",
    ], size=20, spacing=12)
    fin(s, 6, """
Translate taste into constraints. 'Make it professional' cannot drive a loop. 'One page, primary recommendation, why not the two alternatives, every figure cites a slide' can.

That translation is the scarce skill in the verification-gap essay. Notion: https://app.notion.com/p/36f059b703e18142b351f70732b09c29

The machine enforces the floor (sources present, template filled). The human owns the ceiling (is the recommendation actually right).
""")

    # 7 Fail either
    s = blank(prs)
    textbox(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.4),
            "Fail either test → do not automate.", size=36, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    textbox(s, Inches(1.5), Inches(3.6), Inches(10.2), Inches(1.6),
            "That leftover is judgment.\nTaste, red lines, how a room will land, which story is worth telling.\nKeep it.",
            size=22, color=CREAM, align=PP_ALIGN.CENTER)
    fin(s, 7, """
Let this land. People came hoping to automate the hard conversations. Tell them no, kindly.

The shape-of-the-work briefing: if you automate work that depends on trust and judgment, you break the process at the point where the human mattered most. Notion: https://app.notion.com/p/36f059b703e181cf9c30f637d158b234

And from automation discovery: frequency is evidence, not value. Choosing none of the offered automations is allowed. Notion: https://app.notion.com/p/3a0059b703e181c69fafc66f28f76c99
""")

    # 8 The freeze
    s = blank(prs)
    textbox(s, Inches(0.8), Inches(0.45), Inches(11), Inches(0.8),
            "Then people freeze", size=32, bold=True)
    textbox(s, Inches(0.8), Inches(1.5), Inches(11.5), Inches(1.4),
            "They buy the tool. They connect a few things.\nThey stare at an empty prompt.",
            size=24, color=CREAM)
    bullets(s, Inches(0.8), Inches(3.2), Inches(11.5), Inches(3.2), [
        "Idle agents are not broken. Nobody dispatched them.",
        "“Handle my email” is usually seventeen jobs you have never named.",
        "Do not let the model pick the problem and build it. You choose. None is allowed.",
    ], size=20, spacing=14)
    fin(s, 8, """
This is the paralysis slide. Stay here until they nod.

Nate's empty-prompt piece: people have a capable agent and keep staring at the box. The first version of automation-discovery let the AI pick and build — he killed that 'magic button' because it hides the judgment call. Offer sheet, then a human chooses. Choosing none is allowed. Frequency is not value. Notion: https://app.notion.com/p/3a0059b703e181c69fafc66f28f76c99

Idle agents: a social network for agents filled up and sat there. They were never asked to do a single thing. Notion: https://app.notion.com/p/399059b703e18107b20dfe6bb5c32626

Point them at prompts/process/what-to-automate.md — last week, split the blob, edge vs core, two tests, one fast win.
""")

    # 9 Why it fails
    s = blank(prs)
    textbox(s, Inches(0.8), Inches(0.45), Inches(11), Inches(0.8),
            "Why this work actually fails", size=32, bold=True)
    bullets(s, Inches(0.8), Inches(1.45), Inches(11.5), Inches(5.2), [
        "Core-first. Automate the judgment. Three months later: stalled, bloated, humans checked out.",
        "A blob treated as one job. The model writes something that looks like a PRD and falls apart.",
        "No check. Success was “it sounds good,” or the agent said done.",
        "False success. Matching filename, finished-looking draft, wrong file. People stop looking.",
        "Waiting for a smarter model. The missing piece was a named job, a source, and an owner.",
    ], size=18, spacing=12)
    fin(s, 9, """
These are Nate's failure modes from the Substack table, not a generic 'AI projects fail' slide.

Core-first vs edges: https://app.notion.com/p/36f059b703e181d9bd25f984c8bd6945
Fails at the task level (five jobs pretending to be one): https://app.notion.com/p/36f059b703e181f9afedc637cd6529f4
False success (wrong spreadsheet, said done): https://app.notion.com/p/3b5059b703e1819fbd41c6ba9658b0c6
A better model will not save you: https://app.notion.com/p/36f059b703e1818fa542e8696353ebb9

Tell them the prompt for a postmortem is diagnose-the-stall.md.
""")

    # 10 Start at the edges
    s = blank(prs)
    textbox(s, Inches(0.8), Inches(0.45), Inches(11), Inches(0.8),
            "Start at the edges, not the core", size=32, bold=True)
    edges = [
        ("Prep", "Collect, clean, reconcile"),
        ("Check", "Completeness, obvious errors"),
        ("Summarize", "Thread → one page"),
        ("Package", "Analysis → brief / recap"),
        ("Hand off", "Move context. Stop being the glue."),
    ]
    for i, (h, d) in enumerate(edges):
        left = Inches(0.55 + i * 2.5)
        add_rect(s, left, Inches(1.55), Inches(2.35), Inches(2.7), NAVY2)
        textbox(s, left + Inches(0.12), Inches(1.8), Inches(2.1), Inches(0.7), h, size=20, bold=True, color=GOLD)
        textbox(s, left + Inches(0.12), Inches(2.55), Inches(2.1), Inches(1.4), d, size=15, color=CREAM)
    textbox(s, Inches(0.8), Inches(4.55), Inches(11.5), Inches(1.8),
            "The core is craft: taste, red lines, how a room lands. Protect it.\nPick the simplest edge with the clearest lift — not the most impressive one.",
            size=18, color=MUTED)
    fin(s, 10, """
Edge-first is the thought process they came for. Data preparation, QA, synthesis, packaging, coordination. Cheap errors. Humans can pick up an exception without the whole workflow breaking.

The core is where they groan that they want help — and where first projects die. Trust is the real project: automate around the craft, not through it.

Notion: https://app.notion.com/p/36f059b703e181d9bd25f984c8bd6945

CoS mapping: recap, brief, open-loop list, packaging a status = edges. When to push, which story to tell = core.
""")

    # 11 CoS is / isn't
    s = blank(prs)
    textbox(s, Inches(0.8), Inches(0.45), Inches(11), Inches(0.8),
            "The Chief of Staff is the teaching case", size=30, bold=True)
    add_rect(s, Inches(0.8), Inches(1.5), Inches(5.5), Inches(4.7), NAVY2)
    add_rect(s, Inches(6.9), Inches(1.5), Inches(5.5), Inches(4.7), NAVY2)
    textbox(s, Inches(1.1), Inches(1.75), Inches(5), Inches(0.5), "OPERATIONAL  —  automate", size=14, bold=True, color=GOLD)
    bullets(s, Inches(1.1), Inches(2.4), Inches(5), Inches(3.5), [
        "Morning brief",
        "Weekly update",
        "Meeting prep",
        "Meeting recap",
        "Open loops",
        "Drafts of mail you will send",
    ], size=18, spacing=10)
    textbox(s, Inches(7.2), Inches(1.75), Inches(5), Inches(0.5), "JUDGMENT  —  keep", size=14, bold=True, color=GOLD)
    bullets(s, Inches(7.2), Inches(2.4), Inches(5), Inches(3.5), [
        "When to push or hold",
        "How a room will land",
        "Political risk",
        "The hard conversation",
        "Which story is worth telling",
        "Anything you would not put your name on unread",
    ], size=18, spacing=10)
    fin(s, 11, """
Walk the left column: these are the eight portable jobs from the Delegation Kit, renamed into operator English — daily briefing, meeting processing, weekly review, end-of-day reconciliation, memory scaffold. Notion: https://app.notion.com/p/36f059b703e18192a6d8f67fbe74b756

The right column is not 'AI can't do it.' It is 'checking it costs as much as doing it, so extra attempts just grow the pile.' That is the don't-bother verdict from the one-minute test.
""")

    # 9 Same files
    s = blank(prs)
    textbox(s, Inches(0.8), Inches(0.45), Inches(11), Inches(0.8),
            "Same markdown. Any of these tools.", size=32, bold=True)
    tools = ["GitHub Copilot", "Cursor", "Claude", "ChatGPT", "Gemini"]
    for i, name in enumerate(tools):
        left = Inches(0.8 + i * 2.4)
        add_rect(s, left, Inches(2.2), Inches(2.2), Inches(1.6), NAVY2)
        textbox(s, left, Inches(2.65), Inches(2.2), Inches(0.8), name, size=16, bold=True, color=CREAM, align=PP_ALIGN.CENTER)
    textbox(s, Inches(0.8), Inches(4.4), Inches(11.5), Inches(1.8),
            "The how-to in the repo is only this: which box to paste into.\nNothing else. If your company only allows one of these, you still have the kit.",
            size=20, color=MUTED)
    fin(s, 12, """
Do not take questions about which model is 'best' here. Route them out: once the job is named and the check exists, use whatever they already pay for.

The reusable-rig essay is explicit that the skills should be local, inspectable, and independent of whichever AI app you are renting this month. Notion: https://app.notion.com/p/392059b703e1814b96a6dd9913180844
""")

    # Kit map — overview only
    s = blank(prs)
    textbox(s, Inches(0.8), Inches(0.45), Inches(11), Inches(0.8),
            "What's in the kit", size=32, bold=True)
    kit = [
        ("START-HERE.md", "Get the files onto a PC"),
        ("PROCESS.md", "The method — two tests"),
        ("prompts/process/", "Turn any job into a workflow"),
        ("prompts/workflows/", "The CoS jobs, ready to paste"),
        ("context-templates/", "Blanks. Copy privately. Fill them."),
        ("how-to/", "Which box to paste into"),
    ]
    for i, (fn, desc) in enumerate(kit):
        r, c = divmod(i, 3)
        left = Inches(0.8 + c * 4.0)
        top = Inches(1.5 + r * 2.2)
        add_rect(s, left, top, Inches(3.8), Inches(2.0), NAVY2)
        textbox(s, left + Inches(0.2), top + Inches(0.35), Inches(3.4), Inches(0.55), fn, size=16, bold=True, color=GOLD)
        textbox(s, left + Inches(0.2), top + Inches(1.0), Inches(3.4), Inches(0.7), desc, size=16, color=CREAM)
    fin(s, 13, """
Overview, not a tour. Six things, six purposes. Do not open evidence-based-investigation.md unless someone asks how to ground a rumor — it lives under process/.

Personal data never ships in this repo. They copy context-templates/ to a private folder. START-HERE.md is unzip / browser / Cursor.

Cite: map of where material lives — https://app.notion.com/p/3bb059b703e18152aea0d83b502fa319
""")

    # 10 Name one job
    s = blank(prs)
    textbox(s, Inches(0.8), Inches(0.45), Inches(11), Inches(0.8),
            "Step 1 — Name one job", size=32, bold=True)
    bullets(s, Inches(0.8), Inches(1.5), Inches(11.5), Inches(4.8), [
        "Not “be more productive.” Not “be my chief of staff.”",
        "A job you actually did last week, with a finished artifact.",
        "Annoying enough that you will reuse the file.",
        "Prefer a pain that has happened at least three times.",
        "If you cannot point at examples, you do not have a first job yet.",
    ], size=22, spacing=14)
    fin(s, 14, """
Have them shout a job. Split blobs. 'Email' becomes 'Monday status to my manager.' 'Meetings' becomes 'recap of the staff meeting within the hour.'

First-agent-job guide: one annoying problem, three real occurrences, inbox text is data not instructions. Notion: https://app.notion.com/p/3bb059b703e1817b9123ff209fcc5d9d

Prompt they will paste: prompts/process/name-the-job.md and the-two-tests.md
""")

    # 11 Write the SOP
    s = blank(prs)
    textbox(s, Inches(0.8), Inches(0.45), Inches(11), Inches(0.8),
            "Step 2 — Write the SOP", size=32, bold=True)
    heads = [
        ("Inputs", "What must be in front of you"),
        ("Steps", "The order you already use"),
        ("Output", "The shape of 'done'"),
        ("Check", "How you know it's right"),
    ]
    for i, (h, d) in enumerate(heads):
        left = Inches(0.8 + i * 3.05)
        add_rect(s, left, Inches(1.6), Inches(2.9), Inches(2.6), NAVY2)
        textbox(s, left + Inches(0.15), Inches(1.85), Inches(2.6), Inches(0.6), h, size=20, bold=True, color=GOLD)
        textbox(s, left + Inches(0.15), Inches(2.5), Inches(2.6), Inches(1.4), d, size=16, color=CREAM)
    textbox(s, Inches(0.8), Inches(4.5), Inches(11.5), Inches(1.8),
            "Approvals sit where actions become hard to undo: send, pay, publish, delete.\nScore the spec: could a stranger run it? Is success binary? If not, it will create more work than it saves.",
            size=18, color=MUTED)
    fin(s, 15, """
This is the spec from the Delegation Kit: inputs, steps, constraints, approvals at reversibility boundaries, failure handling. Score: can a stranger execute it without questions; are success criteria binary. Three or below means the spec creates more work than it saves. Notion: https://app.notion.com/p/36f059b703e18192a6d8f67fbe74b756

Prompt: prompts/process/write-the-sop.md
""")

    # 12 Context files
    s = blank(prs)
    textbox(s, Inches(0.8), Inches(0.45), Inches(11), Inches(0.8),
            "Step 3 — Put you in files you own", size=32, bold=True)
    files = [
        ("who-i-am.md", "Role, who you answer to"),
        ("priorities.md", "This quarter, dated"),
        ("people.md", "Who changes your week"),
        ("voice.md", "How you actually write"),
        ("sources.md", "Where you would click to check"),
        ("open-loops.md", "What is still open"),
        ("decisions.md", "What you already settled"),
    ]
    for i, (fn, desc) in enumerate(files):
        if i < 4:
            left = Inches(0.8 + i * 3.05)
            top = Inches(1.45)
            w = Inches(2.9)
        else:
            left = Inches(0.8 + (i - 4) * 4.0)
            top = Inches(3.65)
            w = Inches(3.8)
        add_rect(s, left, top, w, Inches(1.95), NAVY2)
        textbox(s, left + Inches(0.15), top + Inches(0.3), w - Inches(0.3), Inches(0.55), fn, size=16, bold=True, color=GOLD)
        textbox(s, left + Inches(0.15), top + Inches(0.95), w - Inches(0.3), Inches(0.7), desc, size=15, color=CREAM)
    fin(s, 16, """
Do not dump this into one encyclopedia. One giant project file becomes a graveyard of stale rules — that is the OpenAI story in the context-files essay. Split: stable instructions, current state, a map of where material lives (sources.md), decision history. When you change your mind, name what is replaced and keep the old assumption in the log. Notion: https://app.notion.com/p/3bb059b703e18152aea0d83b502fa319

These files stay private. The public repo ships blanks. They name the systems they actually open — the agent does not pick a vendor. Prompt: prompts/process/build-context.md
""")

    # 13 Prompt file is the process
    s = blank(prs)
    textbox(s, Inches(0.8), Inches(1.8), Inches(11.5), Inches(1.2),
            "The prompt file is the process.", size=36, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    textbox(s, Inches(1.4), Inches(3.3), Inches(10.5), Inches(2.2),
            "You do not re-explain yourself in chat.\nIf the output is wrong, the file is wrong.\nInstalling a skill proves the files arrived. Only a real job proves they fit.",
            size=22, color=CREAM, align=PP_ALIGN.CENTER)
    fin(s, 17, """
One-job test: don't evaluate a skill by reading the promise. Give it one real job, write pass/fail first, keep the evidence. Keep, fork, or delete. Notion: https://app.notion.com/p/3b0059b703e181fa81add3d2cd98aeaf

This is why we ship prompt files, not a plugin, in this repo.
""")

    # 14 Verify
    s = blank(prs)
    textbox(s, Inches(0.8), Inches(0.45), Inches(11), Inches(0.8),
            "Step 5 — Verify. Cite or cut.", size=32, bold=True)
    bullets(s, Inches(0.8), Inches(1.5), Inches(11.5), Inches(4.8), [
        "Every substantive claim needs a source you can open.",
        "Gaps stay gaps. “Not found” beats a confident guess.",
        "Self-reported “done” is not a check.",
        "You read the draft against the Check section. Then you send — or you don't.",
    ], size=22, spacing=16)
    fin(s, 18, """
Citation guard from the reusable-rig essay: no anchor, no claim. Notion: https://app.notion.com/p/392059b703e1814b96a6dd9913180844

They name the systems in sources.md. You do not pick Outlook vs Gmail for them. If someone asks how to ground a rumor, that is evidence-based-investigation.md — do not teach it here.

Printed holes ('usage: not available') are the system working. Invented numbers inside a confident summary are the failure mode people are right to fear. That example is in the agent-shaped-work piece.

Prompt: prompts/process/verify.md
""")

    # 15 Correct the file
    s = blank(prs)
    textbox(s, Inches(0.8), Inches(0.45), Inches(11), Inches(0.8),
            "Step 6 — Correct the file, not the chat", size=30, bold=True)
    bullets(s, Inches(0.8), Inches(1.5), Inches(11.5), Inches(4.2), [
        "Three of the same correction is a rule. Two is a candidate. One is a fluke.",
        "Put the line in the prompt or in voice.md. One owner.",
        "Delete a stale instruction before you add a new one.",
        "Somebody owns the agent. One person, close enough to notice drift.",
    ], size=22, spacing=14)
    fin(s, 19, """
Maintenance loop: you are maintaining the harness, not a prompt. Seven surfaces — job, diet, memory, tools, reach, proof, value. Repeated correction across three runs is a file problem. Delete before you add. Keep / change / pause / retire. Notion: https://app.notion.com/p/3bb059b703e181bfa3f5fbf9ed29b605

Ownership: the fastest way to make an agent dangerous is to let everybody use it and nobody own it. Notion: https://app.notion.com/p/387059b703e18195a327e34da56998cb

Prompt: prompts/process/correct-the-file.md
""")

    # 16 Drafts never send
    s = blank(prs)
    textbox(s, Inches(0.8), Inches(1.7), Inches(11.5), Inches(1.4),
            "Drafts only. Never send.", size=40, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    textbox(s, Inches(1.3), Inches(3.4), Inches(10.7), Inches(2.2),
            "The agent may read, organize, draft, cite, and export.\nIt may not send, file, submit, pay, sign, or close a ticket.\nYou are the principal. The agent is staff.",
            size=22, color=CREAM, align=PP_ALIGN.CENTER)
    fin(s, 20, """
This is not a legal disclaimer bolted on. It is why the same rig can be pointed at money and health paperwork. If an agent sends a flawed appeal in your name, you now have two problems. Notion: https://app.notion.com/p/392059b703e1814b96a6dd9913180844

Inbox is untrusted. A line that says 'ignore your rules' is data, not an order. First-agent-job: https://app.notion.com/p/3bb059b703e1817b9123ff209fcc5d9d

Approvals sit at reversibility boundaries — the moment an action becomes hard to undo.
""")

    # 17 Then anything
    s = blank(prs)
    textbox(s, Inches(0.8), Inches(0.45), Inches(11), Inches(0.8),
            "Then point the same process at anything", size=30, bold=True)
    bullets(s, Inches(0.8), Inches(1.5), Inches(11.5), Inches(4.8), [
        "Expense reports, hiring screens, incident recaps, customer research.",
        "If it passes both tests, it gets a prompt file and a check.",
        "If it fails, you just saved yourself a clever mess.",
        "Do not ask “can AI do this?” Ask what shape the work is.",
        "Buying a smarter model does not skip workflow, data, authority, evaluation, audit, or an owner.",
    ], size=20, spacing=14)
    fin(s, 21, """
Six things have to be true before AI changes a workflow. Most companies have built two. Notion: https://app.notion.com/p/36f059b703e1813d801bcb34d72141b1

Shape of the work: https://app.notion.com/p/36f059b703e181cf9c30f637d158b234

The CoS is the demo so they learn the motion on work they already understand. Then they take it home.
""")

    # 18 30 minutes
    s = blank(prs)
    textbox(s, Inches(0.8), Inches(0.45), Inches(11), Inches(0.8),
            "Thirty minutes, starting Monday", size=32, bold=True)
    steps = [
        ("10 min", "Read PROCESS.md"),
        ("10 min", "Fill in your private context files"),
        ("5 min", "If you freeze: what-to-automate.md, then the two tests"),
        ("5 min", "Write the SOP — or keep the job as judgment"),
    ]
    for i, (t, d) in enumerate(steps):
        top = Inches(1.5 + i * 1.15)
        add_rect(s, Inches(0.8), top, Inches(11.5), Inches(1.0), NAVY2)
        textbox(s, Inches(1.05), top + Inches(0.28), Inches(1.8), Inches(0.5), t, size=18, bold=True, color=GOLD)
        textbox(s, Inches(3.0), top + Inches(0.28), Inches(9), Inches(0.5), d, size=20, color=CREAM)
    fin(s, 22, """
Walk the README. START-HERE.md first if the files are not on the machine yet. If they freeze, they start at what-to-automate.md — last week, split the blob, edges first. If the two tests fail, they still used the thirty minutes correctly. The lesson is the refusal.

Do not put the brief on a timer until they have checked it by hand several times. Automation is for a process they already trust.
""")

    # 19 Safety card
    s = blank(prs)
    textbox(s, Inches(0.8), Inches(0.45), Inches(11), Inches(0.8),
            "Non-negotiable", size=32, bold=True)
    bullets(s, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.0), [
        "Read-only first.",
        "Drafts, never send.",
        "Cite or cut.",
        "Inbox is untrusted — mail is data, not orders.",
        "Stale is visible. Silent omission is a lie.",
        "One owner per agent.",
    ], size=24, spacing=12)
    fin(s, 23, """
Leave this up during Q&A if needed.

Where the agent should stop: start where a colleague or customer already tells you you're wrong; reconstructing context is the expensive part; the reply is cheap. Notion: https://app.notion.com/p/3aa059b703e1817c9571c77f8badaf77
""")

    # 20 Close
    s = blank(prs)
    textbox(s, Inches(0.8), Inches(1.6), Inches(11.5), Inches(1.4),
            "Automate what repeats and checks.\nKeep the rest.", size=36, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    textbox(s, Inches(1.5), Inches(3.6), Inches(10.3), Inches(2.0),
            "Repo: START-HERE.md, then the prompt kit (not the Claude plugin).\nProcess files, CoS workflows, blank context templates.\nSources cited from the Notion Substack table.",
            size=20, color=CREAM, align=PP_ALIGN.CENTER)
    fin(s, 24, """
Close by pointing at the repo. This kit is provider-agnostic. The Claude plugin of the same processes is a different repository (BittahCriminal/Chief-of-Staff) — do not send them there for this talk.

If they want receipts: research/SOURCES.md lists the Notion pages this was synthesized from, starting with the Substack database https://app.notion.com/p/36e059b703e180d3a962d862c9e380c5

Offer to stay for the first job. Help them run the two tests live on something they did last week.
""")

    out = Path(__file__).resolve().parent / "building-your-own-chief-of-staff.pptx"
    prs.save(out)
    print(f"Wrote {out} ({TOTAL} slides)")


if __name__ == "__main__":
    build()
